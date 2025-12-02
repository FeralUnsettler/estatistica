# app.py — REPARA v3.0
# Recursos principais: multi-upload, clustering, modelo de turnover, summaries, export PDF/PPTX
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import plotly.express as px
from wordcloud import WordCloud
import io
import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

# ML
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans
from sklearn.model_selection import train_test_split
from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
from sklearn.metrics import accuracy_score, roc_auc_score, mean_squared_error, classification_report
import joblib
from sklearn.feature_extraction.text import TfidfVectorizer

st.set_page_config(
    page_title="REPARA • Plataforma 3.0",
    layout="wide",
)

st.title("🚀 REPARA — Plataforma Analítica v3.0")
st.write("Multi-upload • Clustering • Previsão de turnover • Resumos • Export PDF/PPTX")

# ---------------------------
# Helpers
# ---------------------------
def read_csv_any(uploaded_file):
    try:
        return pd.read_csv(uploaded_file, engine="python", sep=None)
    except Exception:
        uploaded_file.seek(0)
        return pd.read_csv(uploaded_file)

def infer_columns(df):
    """Detect common columns and return a map of names found (best-effort)."""
    cols = {c.lower(): c for c in df.columns}
    mapping = {}
    # candidates
    for target in ['age','gender','race','main_pain','skills','tenure','left','turnover','hired','exit','experience']:
        if target in cols:
            mapping[target] = cols[target]
    # company
    for target in ['company_size','hr_challenges','industry','revenue']:
        if target in cols:
            mapping[target] = cols[target]
    return mapping

def extract_key_phrases(texts, topk=5):
    # Extractive approach using TF-IDF: return topk terms (unigrams/bigrams)
    if not isinstance(texts, (list, pd.Series)):
        return []
    texts = [str(t) for t in texts if pd.notna(t)]
    if not texts:
        return []
    vec = TfidfVectorizer(ngram_range=(1,2), max_features=200, stop_words='english')
    X = vec.fit_transform(texts)
    sums = X.sum(axis=0)
    terms = vec.get_feature_names_out()
    scores = [(terms[i], sums[0,i]) for i in range(len(terms))]
    scores = sorted(scores, key=lambda x: x[1], reverse=True)
    return [t for t,s in scores[:topk]]

def make_wordcloud(text):
    wc = WordCloud(width=900, height=400, background_color="white").generate(text)
    fig, ax = plt.subplots(figsize=(12,4))
    ax.imshow(wc, interpolation='bilinear')
    ax.axis('off')
    return fig

def try_to_create_target(df, mapping):
    # Tries to infer a binary target 'left_90' or 'left' or create proxy by tenure < threshold
    if 'left' in mapping:
        target_col = mapping['left']
        # Normalize to binary
        if df[target_col].dtype == object:
            y = df[target_col].str.lower().isin(['yes','y','1','true','sim','s']).astype(int)
        else:
            y = (df[target_col] != 0).astype(int)
        return y, target_col
    if 'turnover' in mapping:
        t = mapping['turnover']
        y = df[t].apply(lambda x: 1 if str(x).strip().lower() in ['yes','y','sim','1','true'] else 0)
        return y, t
    # fallback: if tenure exists, create proxy: tenure < 6 months -> likely left (1)
    if 'tenure' in mapping:
        tenure_col = mapping['tenure']
        try:
            y = (pd.to_numeric(df[tenure_col], errors='coerce') < 6).astype(int)
            return y, tenure_col
        except Exception:
            return None, None
    return None, None

# ---------------------------
# Sidebar: uploads
# ---------------------------
st.sidebar.header("Uploads & Config")
st.sidebar.write("Envie 1 ou mais CSVs. O app tenta interpretar automaticamente.")

candidate_uploads = st.sidebar.file_uploader("CSV(s) de Candidatos — (pode enviar múltiplos)", type=['csv'], accept_multiple_files=True)
company_uploads = st.sidebar.file_uploader("CSV(s) de Empresas — (pode enviar múltiplos)", type=['csv'], accept_multiple_files=True)

# advanced params
n_clusters = st.sidebar.number_input("K para KMeans (clusters)", min_value=2, max_value=12, value=4)
random_seed = st.sidebar.number_input("Random seed", min_value=0, value=42)
perform_ml = st.sidebar.checkbox("Habilitar ML (clustering + turnover model)", value=True)
save_model_button = st.sidebar.checkbox("Salvar modelo treinado (joblib) após treino", value=True)

# ---------------------------
# Load datasets
# ---------------------------
@st.cache_data
def load_concat(files):
    dfs = []
    for f in files:
        try:
            df = read_csv_any(f)
            dfs.append(df)
        except Exception as e:
            st.warning(f"Erro ao ler {getattr(f,'name',str(f))}: {e}")
    if dfs:
        return pd.concat(dfs, ignore_index=True)
    else:
        return pd.DataFrame()

candidates_df = load_concat(candidate_uploads) if candidate_uploads else pd.DataFrame()
companies_df = load_concat(company_uploads) if company_uploads else pd.DataFrame()

# ---------------------------
# Tabs
# ---------------------------
tabs = st.tabs(["Overview", "Candidatos", "Empresas", "Clustering", "Turnover Model", "Export"])
# ---------------------------
# Overview
# ---------------------------
with tabs[0]:
    st.header("📋 Overview")
    st.write("Resumo rápido dos datasets carregados.")
    col1, col2 = st.columns(2)
    col1.metric("Candidatos (lin.)", len(candidates_df))
    col2.metric("Empresas (lin.)", len(companies_df))
    if not candidates_df.empty:
        st.subheader("Colunas — Candidatos")
        st.write(list(candidates_df.columns))
    if not companies_df.empty:
        st.subheader("Colunas — Empresas")
        st.write(list(companies_df.columns))

# ---------------------------
# Candidatos tab
# ---------------------------
with tabs[1]:
    st.header("👤 Candidatos — Exploração")
    if candidates_df.empty:
        st.info("Faça upload dos CSVs de candidatos no menu lateral.")
    else:
        st.subheader("Preview")
        st.dataframe(candidates_df.head(200))

        mapping = infer_columns(candidates_df)
        st.subheader("Colunas detectadas (mapeamento)")
        st.json(mapping)

        # demographics
        st.markdown("### Demografia rápida")
        demo_cols = []
        for k in ['age','gender','race','experience']:
            if k in mapping:
                demo_cols.append(mapping[k])
        if demo_cols:
            for c in demo_cols:
                fig = px.histogram(candidates_df, x=c, title=f"Distribuição — {c}")
                st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("Nenhuma coluna demográfica padrão detectada.")

        # free text wordcloud and keyphrases
        if 'main_pain' in mapping:
            st.markdown("### Textual: Wordcloud & Key-phrases")
            text = " ".join(candidates_df[mapping['main_pain']].dropna().astype(str).tolist())
            st.pyplot(make_wordcloud(text))
            kps = extract_key_phrases(candidates_df[mapping['main_pain']].dropna().astype(str), topk=8)
            st.write("Key phrases:", kps)
        else:
            st.info("Coluna de texto com 'main_pain' não detectada; para resumos, nomeie a coluna como 'main_pain' ou similar.")

        # Save cleaned CSV
        st.markdown("### Download dataset (limpo)")
        buffer = io.StringIO()
        candidates_df.to_csv(buffer, index=False)
        st.download_button("📥 Baixar candidatos.csv (processado)", buffer.getvalue(), "candidatos_processed.csv", mime="text/csv")


# ---------------------------
# Empresas tab
# ---------------------------
with tabs[2]:
    st.header("🏢 Empresas — Exploração")
    if companies_df.empty:
        st.info("Faça upload dos CSVs de empresas no menu lateral.")
    else:
        st.subheader("Preview")
        st.dataframe(companies_df.head(200))

        mapping = infer_columns(companies_df)
        st.subheader("Colunas detectadas (mapeamento)")
        st.json(mapping)

        if 'hr_challenges' in mapping:
            st.markdown("### Top desafios de RH")
            top = companies_df[mapping['hr_challenges']].dropna().astype(str).value_counts().head(10)
            fig = px.bar(x=top.values, y=top.index, orientation='h', title="Top desafios de RH")
            st.plotly_chart(fig, use_container_width=True)

        st.markdown("### Download dataset (limpo)")
        buffer = io.StringIO()
        companies_df.to_csv(buffer, index=False)
        st.download_button("📥 Baixar empresas.csv (processado)", buffer.getvalue(), "empresas_processed.csv", mime="text/csv")

# ---------------------------
# Clustering tab
# ---------------------------
with tabs[3]:
    st.header("🔍 Clustering — Agrupe candidatos por perfil")
    if candidates_df.empty:
        st.info("Faça upload dos dados de candidatos para usar clustering.")
    else:
        # select numeric columns for clustering (or skills vectorize)
        numeric_cols = candidates_df.select_dtypes(include=[np.number]).columns.tolist()
        # if numeric not enough, try to encode age or experience
        if not numeric_cols:
            st.warning("Sem colunas numéricas suficientes; tentando extrair 'age' se existir.")
            mapping = infer_columns(candidates_df)
            if 'age' in mapping:
                candidates_df['_age_num'] = pd.to_numeric(candidates_df[mapping['age']], errors='coerce').fillna(-1)
                numeric_cols = ['_age_num']

        if not numeric_cols:
            st.error("Não há features numéricas para clustering. Considere adicionar colunas quantitativas ou usar skills numéricas.")
        else:
            st.write("Features usadas para clustering:", numeric_cols)
            X = candidates_df[numeric_cols].fillna(-1).values
            scaler = StandardScaler()
            Xs = scaler.fit_transform(X)
            kmeans = KMeans(n_clusters=n_clusters, random_state=random_seed)
            labels = kmeans.fit_predict(Xs)
            candidates_df['cluster'] = labels
            st.success(f"Clusters gerados: {n_clusters}")
            fig = px.scatter_matrix(pd.DataFrame(Xs, columns=numeric_cols).assign(cluster=labels),
                                    dimensions=numeric_cols[:3] if len(numeric_cols)>=3 else numeric_cols,
                                    color=labels.astype(str),
                                    title="Scatter matrix (features escaladas)")
            st.plotly_chart(fig, use_container_width=True)
            st.dataframe(candidates_df[[*numeric_cols, 'cluster']].head(200))

            # cluster counts
            st.markdown("### Tamanho por cluster")
            counts = candidates_df['cluster'].value_counts().sort_index()
            fig2 = px.bar(x=counts.index.astype(str), y=counts.values, labels={'x':'cluster','y':'count'})
            st.plotly_chart(fig2, use_container_width=True)

            # download labels
            out_buf = io.StringIO()
            candidates_df.to_csv(out_buf, index=False)
            st.download_button("📥 Baixar candidatos com cluster", out_buf.getvalue(), "candidatos_clustered.csv", mime="text/csv")

# ---------------------------
# Turnover Model tab
# ---------------------------
with tabs[4]:
    st.header("📈 Turnover Prediction (modelo simples)")
    if candidates_df.empty:
        st.info("Carregue CSVs de candidatos para treinar o modelo.")
    else:
        mapping = infer_columns(candidates_df)
        y, target_col = try_to_create_target(candidates_df, mapping)

        if y is None or len(y.unique()) <= 1:
            st.warning("Não foi possível identificar um target binário de 'left' / 'turnover'. Se existir uma coluna com saída, nomeie-a 'left' ou 'turnover'.")
            st.info("Como fallback, você pode criar um proxy: por exemplo, tenure < 6 meses -> left.")
            if st.button("Criar proxy 'left' baseado em tenure < 6 meses"):
                if 'tenure' in mapping:
                    candidates_df['left_proxy'] = (pd.to_numeric(candidates_df[mapping['tenure']], errors='coerce') < 6).astype(int)
                    y = candidates_df['left_proxy']
                    target_col = 'left_proxy'
                    st.success("Proxy criado como left_proxy.")
                else:
                    st.error("Não há coluna 'tenure' detectada.")
        if y is not None and len(y.unique())>1:
            st.write("Target identificado (col):", target_col)
            # features: try numeric + simple encoding of gender/race
            X = candidates_df.select_dtypes(include=[np.number]).copy()
            if X.shape[1] < 2:
                # try to encode categorical small-cardinality columns
                for col in ['gender','race','experience']:
                    if col in mapping:
                        X[col] = pd.factorize(candidates_df[mapping[col]])[0]
            # drop rows without target
            df_ml = pd.concat([X, y], axis=1).dropna()
            if df_ml.shape[0] < 20:
                st.warning("Poucos registros com target e features válidas para treinar (recomendado >= 20).")
            else:
                Xf = df_ml.drop(columns=[target_col]) if target_col in df_ml.columns else df_ml.drop(columns=[y.name])
                yf = df_ml[y.name] if hasattr(y,'name') and y.name in df_ml.columns else df_ml.iloc[:,-1]
                X_train, X_test, y_train, y_test = train_test_split(Xf, yf, test_size=0.25, random_state=random_seed)
                # choose model type
                model = RandomForestClassifier(n_estimators=100, random_state=random_seed)
                model.fit(X_train, y_train)
                preds = model.predict(X_test)
                probs = model.predict_proba(X_test)[:,1] if hasattr(model, "predict_proba") else None

                st.subheader("Resultado do treino")
                st.write("Treino:", X_train.shape, "Teste:", X_test.shape)
                st.write("Accuracy:", accuracy_score(y_test, preds))
                if probs is not None:
                    try:
                        st.write("ROC AUC:", roc_auc_score(y_test, probs))
                    except Exception:
                        pass
                st.text("Classificação (test set):")
                st.text(classification_report(y_test, preds))

                # Feature importances
                fi = pd.Series(model.feature_importances_, index=Xf.columns).sort_values(ascending=False).head(15)
                st.subheader("Importância das features")
                st.write(fi)
                fig = px.bar(x=fi.index, y=fi.values, labels={'x':'feature','y':'importance'}, title="Feature importances")
                st.plotly_chart(fig, use_container_width=True)

                # Save model
                if save_model_button:
                    model_path = f"repara_model_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.joblib"
                    joblib.dump(model, model_path)
                    st.success(f"Modelo salvo em {model_path}")
                    with open(model_path, "rb") as f:
                        st.download_button("📥 Baixar modelo (.joblib)", f, file_name=os.path.basename(model_path))

# ---------------------------
# Export tab (PDF + PPTX)
# ---------------------------
with tabs[5]:
    st.header("📦 Export — Relatório e Pitch (PDF / PPTX)")
    st.write("Gera um relatório executivo e um PPTX com slides baseados nos dados carregados.")

    if st.button("Gerar relatório PDF + PPTX"):
        # Basic PDF using fpdf
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt="REPARA — Relatório Executivo", ln=True)
        pdf.ln(6)
        pdf.cell(200, 8, txt=f"Candidatos: {len(candidates_df)} linhas", ln=True)
        pdf.cell(200, 8, txt=f"Empresas: {len(companies_df)} linhas", ln=True)

        # add top pains and top hr challenges
        pdf.ln(6)
        if not candidates_df.empty and 'main_pain' in infer_columns(candidates_df):
            mp_col = infer_columns(candidates_df)['main_pain']
            mp_top = candidates_df[mp_col].dropna().astype(str).value_counts().head(5)
            pdf.cell(200, 8, txt="Top dores (candidatos):", ln=True)
            for t,v in mp_top.items():
                pdf.cell(200,6, txt=f"- {t} ({v})", ln=True)
        if not companies_df.empty and 'hr_challenges' in infer_columns(companies_df):
            hr_col = infer_columns(companies_df)['hr_challenges']
            hr_top = companies_df[hr_col].dropna().astype(str).value_counts().head(5)
            pdf.ln(4)
            pdf.cell(200, 8, txt="Top RH challenges (empresas):", ln=True)
            for t,v in hr_top.items():
                pdf.cell(200,6, txt=f"- {t} ({v})", ln=True)

        pdf_output = pdf.output(dest="S").encode("latin-1")
        st.download_button("📥 Baixar relatório (PDF)", data=pdf_output, file_name="repara_report.pdf", mime="application/pdf")

        # PPTX creation
        prs = Presentation()
        # Title
        s0 = prs.slides.add_slide(prs.slide_layouts[0])
        s0.shapes.title.text = "REPARA — Relatório Executivo"
        s0.placeholders[1].text = f"Candidatos: {len(candidates_df)}  •  Empresas: {len(companies_df)}"

        # Problem slide
        s1 = prs.slides.add_slide(prs.slide_layouts[1])
        s1.shapes.title.text = "Problema"
        tf = s1.shapes.placeholders[1].text_frame
        tf.text = "Candidatos afirmativos relatam invisibilidade; empresas relatam turnover e onboarding ineficiente."

        # Data slides: add images temporarily from Matplotlib/Plotly if possible
        # Example: top pains as text slide
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Top dores (candidatos)"
        tf = s2.shapes.placeholders[1].text_frame
        if not candidates_df.empty and 'main_pain' in infer_columns(candidates_df):
            mp_col = infer_columns(candidates_df)['main_pain']
            mp_top = candidates_df[mp_col].dropna().astype(str).value_counts().head(5)
            for t,v in mp_top.items():
                tf.add_paragraph().text = f"- {t} ({v})"
        else:
            tf.text = "Dados não disponíveis."

        # Save pptx to buffer
        pptx_path = "repara_report.pptx"
        prs.save(pptx_path)
        with open(pptx_path, "rb") as f:
            st.download_button("📥 Baixar PPTX", f, file_name=pptx_path)

st.write("— fim do app v3.0 —")
   
# End of app.py
